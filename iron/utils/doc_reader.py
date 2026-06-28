"""二进制文档解析 — 支持 docx/pdf/excel/pptx 等格式

用途：read_file 工具检测到二进制文件时，根据扩展名调用对应库提取文本，
避免 AI 反复用 run_command 试错（每次命令不同还会触发重复授权询问）。

设计原则：
- 懒加载：只有真正需要读取时才 import 对应库，避免启动时加载未安装的依赖
- 优雅降级：库未安装时返回清晰的安装提示（一次性），不抛异常
- 统一返回 (content, language, error) 三元组，调用方统一处理
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional


# 扩展名 → 解析函数 的映射（懒加载）
def read_docx(file_path: Path) -> tuple[str, str, Optional[str]]:
    """读取 .docx 文档，提取段落 + 表格内容

    Returns: (content, language, error)
        content: 提取的文本内容
        language: "markdown"
        error: 错误信息（成功时为 None）
    """
    try:
        import docx  # python-docx
    except ImportError:
        return "", "text", "需要安装 python-docx：pip install python-docx"

    try:
        doc = docx.Document(str(file_path))
        parts: list[str] = []

        # 段落
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 保留标题层级（Heading 1/2/3 → markdown #）
                style = (para.style.name or "").lower()
                if "heading 1" in style:
                    parts.append(f"# {text}")
                elif "heading 2" in style:
                    parts.append(f"## {text}")
                elif "heading 3" in style:
                    parts.append(f"### {text}")
                else:
                    parts.append(text)

        # 表格
        for i, table in enumerate(doc.tables):
            parts.append(f"\n[表格 {i + 1}]")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))

        content = "\n".join(parts)
        return content, "markdown", None
    except Exception as e:
        return "", "text", f"读取 docx 失败: {e}"


def read_pdf(file_path: Path) -> tuple[str, str, Optional[str]]:
    """读取 .pdf 文档，提取文本

    优先用 pdfplumber（表格支持好），不可用时用 PyPDF2。
    """
    # 优先 pdfplumber
    try:
        import pdfplumber
        try:
            parts: list[str] = []
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        parts.append(f"--- 第 {i + 1} 页 ---\n{text}")
                    # 表格
                    tables = page.extract_tables()
                    for j, table in enumerate(tables):
                        parts.append(f"\n[表格 {j + 1}]")
                        for row in table:
                            cells = [c or "" for c in row]
                            parts.append(" | ".join(cells))
            content = "\n\n".join(parts)
            return content, "markdown", None
        except Exception as e:
            return "", "text", f"读取 pdf 失败: {e}"
    except ImportError:
        pass

    # fallback PyPDF2
    try:
        import PyPDF2
        try:
            parts = []
            with open(str(file_path), "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        parts.append(f"--- 第 {i + 1} 页 ---\n{text}")
            content = "\n\n".join(parts)
            return content, "markdown", None
        except Exception as e:
            return "", "text", f"读取 pdf 失败: {e}"
    except ImportError:
        return "", "text", "需要安装 pdfplumber 或 PyPDF2：pip install pdfplumber"


def read_excel(file_path: Path) -> tuple[str, str, Optional[str]]:
    """读取 .xlsx/.xls 电子表格"""
    try:
        import openpyxl
    except ImportError:
        return "", "text", "需要安装 openpyxl：pip install openpyxl"

    try:
        # .xls（旧格式）openpyxl 不支持，需要 xlrd
        if file_path.suffix.lower() == ".xls":
            try:
                import xlrd
            except ImportError:
                return "", "text", "读取 .xls 需要安装 xlrd：pip install xlrd"
            try:
                book = xlrd.open_workbook(str(file_path))
                parts: list[str] = []
                for sheet in book.sheets():
                    parts.append(f"### 工作表: {sheet.name}")
                    for r in range(sheet.nrows):
                        row = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                        parts.append(" | ".join(row))
                return "\n".join(parts), "markdown", None
            except Exception as e:
                return "", "text", f"读取 xls 失败: {e}"

        # .xlsx
        wb = openpyxl.load_workbook(str(file_path), data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"### 工作表: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                if any(c is not None for c in row):
                    cells = [str(c) if c is not None else "" for c in row]
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts), "markdown", None
    except Exception as e:
        return "", "text", f"读取 excel 失败: {e}"


def read_pptx(file_path: Path) -> tuple[str, str, Optional[str]]:
    """读取 .pptx 演示文稿，提取每页的文本框内容"""
    try:
        import pptx  # python-pptx
    except ImportError:
        return "", "text", "需要安装 python-pptx：pip install python-pptx"

    try:
        prs = pptx.Presentation(str(file_path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"--- 第 {i + 1} 页 ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts), "markdown", None
    except Exception as e:
        return "", "text", f"读取 pptx 失败: {e}"


# 扩展名 → 解析函数 映射
_DOC_READERS: dict[str, callable] = {
    ".docx": read_docx,
    ".pdf": read_pdf,
    ".xlsx": read_excel,
    ".xls": read_excel,
    ".pptx": read_pptx,
}


def is_supported_doc(file_path: Path) -> bool:
    """判断文件是否是内置支持的文档格式"""
    return file_path.suffix.lower() in _DOC_READERS


def read_document(file_path: Path) -> tuple[str, str, Optional[str]]:
    """读取二进制文档，根据扩展名调用对应解析器

    Returns: (content, language, error)
        content: 提取的文本内容（失败时为空）
        language: "markdown"（成功）或 "text"（失败）
        error: 错误信息（成功时为 None）；库未安装时返回安装提示
    """
    ext = file_path.suffix.lower()
    reader = _DOC_READERS.get(ext)
    if reader is None:
        return "", "text", f"不支持的文档格式: {ext}"
    return reader(file_path)
